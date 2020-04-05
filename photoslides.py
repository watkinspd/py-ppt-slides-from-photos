#!/usr/bin/env python3

from pptx import Presentation
from pptx.util import Inches
import os
import sys
import getopt

def main(argv):
    pathtofiles = '.'
    outputfile = 'test.pptx'
    templatefile = './template.pptx'
    sort = 'm'
    try:
        opts, args = getopt.getopt(argv,"hfi:o:t:",["imagelocationdirectory=",
                                                "outputfile=",
                                                "templatefile=",
                                                "sortbyfilename="])
    except getopt.GetoptError:
        print ('try python photoslides.py -h ')
        sys.exit(2)
    for opt, arg in opts:
        if opt == '-h':
            print ('photoslides [options]\n'
            'The following options can be specified:\n'
            '-i, --imagelocationdirectory=<dirOfImages> Must be absolute reference to the directory where the images are located. Do not unclude "/" at end. Default is the current directory.\n'
            '-o, --outputfile=<output.pptx> Name of powerpoint file to generate. Defaults to test.pptx\n'
            '-t, --templatefile=<template.pptx> Source presentation to use. Defaults to ./template.pptx This tool adds the blank slide style which is assumed to'
            ' be 7th in sequence of the master layouts.\n'
            '-f, --sortbyfilename Sort image files by filename. Default sort is by date modified.\n')
            sys.exit()
        elif opt in ("-i", "--imagelocationdirectory"):
            pathtofiles = arg
        elif opt in ("-o", "--outputfile"):
            outputfile = arg
        elif opt in ("-t", "--templatefile"):
            templatefile = arg
        elif opt in ("-f", "--sortbyfilename"):
            sort = 'f'

    prs = Presentation(templatefile)
 
    blank_slide_layout = prs.slide_layouts[6]

    left = top = Inches(0)
    height = prs.slide_height

    folder = os.fsencode(pathtofiles)

    filenames = []

    for file in os.listdir(folder):
        filename = os.fsdecode(file)
        if filename.endswith( ('.jpeg', '.png', '.gif', '.tiff') ): # image files supported by powerpint...
            filenames.append(pathtofiles + "/" + filename)

    if sort == "f":
        filenames.sort()
    else: 
        filenames.sort(key=os.path.getmtime)

    for file in filenames:
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(file, left, top, height=height)

    prs.save(outputfile)

if __name__ == "__main__":
   main(sys.argv[1:])